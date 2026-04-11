"""
Generate PowerPoint presentation for ECR 2.0 Pipeline.
Output: presentasi_ecr2.pptx

Slides:
1.  Title Slide
2.  Daftar Isi
3.  Latar Belakang & Analisis Masalah
4.  Rumusan Masalah
5.  Tujuan Penelitian
6.  Solusi: Overview Pipeline ECR 2.0
7.  Tahap 1: Preprocessing & Sentiment Analysis
8.  Tahap 2: Konstruksi Graf Interaksi
9.  Tahap 3: User Leaning & Community Detection
10. Tahap 4: ECR 2.0 & Validasi Statistik
11. Tahap 5: Analisis Lanjutan
12. Implementasi: Web GUI (NiceGUI)
13. Dataset & Eksperimen
14. Hasil: ECR — Leiden vs Infomap
15. Hasil: Community Quality Metrics
16. Hasil: Visualisasi (1) - Sentimen & Leaning
17. Hasil: Visualisasi (2) - Komunitas & ECR
18. Hasil: Visualisasi (3) - Diffusion & Network
19. Analisis Akhir
20. Kesimpulan
21. Saran & Future Work
22. Terima Kasih
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = Path(__file__).parent
VIZ_DIR = OUT_DIR / "hasil_boikot" / "700K Data"
VIZ_DIR2 = OUT_DIR / "pipeline_out" / "visualizations"

# Colors
DARK_BLUE = RGBColor(0x1a, 0x23, 0x7e)
MED_BLUE = RGBColor(0x28, 0x35, 0x93)
LIGHT_BLUE = RGBColor(0x42, 0xa5, 0xf5)
ACCENT_BLUE = RGBColor(0xe3, 0xf2, 0xfd)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
GREEN = RGBColor(0x2e, 0x7d, 0x32)
RED = RGBColor(0xc6, 0x28, 0x28)
ORANGE = RGBColor(0xef, 0x6c, 0x00)


def find_viz(name):
    for d in [VIZ_DIR, VIZ_DIR2]:
        p = d / name
        if p.exists():
            return str(p)
    return None


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, left, top, width, height,
                     font_size=16, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox


def make_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark blue title bar at top of slide."""
    bar = add_shape_bg(slide, Inches(0), Inches(0),
                       Inches(13.333), Inches(1.2), DARK_BLUE)

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                title_text, font_size=28, bold=True, color=WHITE)

    if subtitle_text:
        add_textbox(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=RGBColor(0xbb, 0xde, 0xfb))


def make_two_col_header(slide, title):
    """Title bar + return y position for content."""
    make_title_bar(slide, title)
    return Inches(1.4)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                "Pengukuran Echo Chamber di Media Sosial Indonesia\n"
                "Menggunakan Analisis Sentimen dan Deteksi Komunitas\n"
                "dengan Framework Echo Chamber Ratio (ECR) 2.0",
                font_size=32, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
                "Amalia Putri",
                font_size=22, bold=True, color=LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
                "Teknik Informatika\n"
                "Sekolah Teknik Elektro dan Informatika\n"
                "Institut Teknologi Bandung (ITB)",
                font_size=16, color=RGBColor(0xbb, 0xde, 0xfb),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
                "2025",
                font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 2: Daftar Isi
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Daftar Isi")

    items = [
        "1.  Latar Belakang & Analisis Masalah",
        "2.  Rumusan Masalah & Tujuan",
        "3.  Solusi: Pipeline ECR 2.0",
        "4.  Tahapan Pipeline (5 Tahap Utama)",
        "5.  Implementasi: Web GUI",
        "6.  Dataset & Eksperimen",
        "7.  Hasil & Analisis",
        "8.  Kesimpulan & Saran",
    ]
    add_bullet_slide(slide, items,
                     Inches(1.5), Inches(1.8), Inches(10), Inches(5),
                     font_size=20, color=BLACK, spacing=Pt(12))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 3: Latar Belakang
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Latar Belakang & Analisis Masalah")

    left_items = [
        "Fenomena Echo Chamber:",
        "• 73% populasi Indonesia aktif di media sosial",
        "• Algoritma rekomendasi memperkuat filter bubble",
        "• Pengguna hanya terpapar informasi sejalan pandangannya",
        "• Memperkuat polarisasi & mempercepat penyebaran misinformasi",
        "",
        "Keterbatasan Pendekatan Existing:",
        "• Network-based: hanya analisis struktur graf,",
        "  mengabaikan konten dan sentimen",
        "• Content-based: capture sentimen tapi miss",
        "  hubungan interaksi antar pengguna",
        "• Belum ada pipeline otomatis end-to-end",
        "  untuk konteks bahasa Indonesia",
    ]
    add_bullet_slide(slide, left_items,
                     Inches(0.5), Inches(1.5), Inches(6), Inches(5.5),
                     font_size=15, color=BLACK)

    right_items = [
        "Gap yang Diidentifikasi:",
        "",
        "• Studi echo chamber di Indonesia masih",
        "  menggunakan pendekatan manual",
        "• Tidak ada validasi statistik yang kuat",
        "  untuk hasil pengukuran echo chamber",
        "• Model sentimen generic tidak optimal",
        "  untuk bahasa Indonesia informal",
        "• Tidak ada perbandingan antar algoritma",
        "  deteksi komunitas",
        "",
        "→ Diperlukan pipeline otomatis yang",
        "  menggabungkan analisis sentimen,",
        "  deteksi komunitas, dan validasi statistik",
    ]
    add_bullet_slide(slide, right_items,
                     Inches(6.8), Inches(1.5), Inches(6), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 4: Rumusan Masalah
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Rumusan Masalah & Tujuan Penelitian")

    items_rm = [
        "Rumusan Masalah:",
        "",
        "1. Bagaimana mengukur echo chamber secara kuantitatif di",
        "   percakapan media sosial Indonesia?",
        "",
        "2. Bagaimana memvalidasi hasil pengukuran echo chamber",
        "   secara statistik agar hasilnya reliable?",
        "",
        "3. Bagaimana membangun pipeline otomatis yang dapat",
        "   digunakan oleh peneliti tanpa memerlukan coding?",
    ]
    add_bullet_slide(slide, items_rm,
                     Inches(0.5), Inches(1.5), Inches(6), Inches(5),
                     font_size=16, color=BLACK)

    items_tj = [
        "Tujuan Penelitian:",
        "",
        "1. Mengembangkan pipeline otomatis end-to-end untuk",
        "   mengukur echo chamber menggunakan ECR 2.0",
        "",
        "2. Mengimplementasikan multiple validation metrics:",
        "   Bootstrap CI, Permutation Test, Community Quality",
        "",
        "3. Membangun web GUI interaktif sehingga pipeline",
        "   dapat digunakan tanpa pengetahuan programming",
        "",
        "4. Menguji pipeline pada dataset nyata percakapan",
        "   di media sosial Indonesia",
    ]
    add_bullet_slide(slide, items_tj,
                     Inches(6.8), Inches(1.5), Inches(6), Inches(5),
                     font_size=16, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 5: Solusi Overview
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Solusi: Pipeline ECR 2.0 — Overview")

    # Pipeline steps as boxes
    steps = [
        ("1", "Input CSV", "Dataset Twitter/X\n(from_id, username, text, reply, retweet, mention)"),
        ("2", "Sentiment Analysis", "Fine-tuned XLM-RoBERTa\n→ p_neg, p_neu, p_pos per tweet"),
        ("3", "Graf Interaksi", "Directed weighted graph\nNode = username, Edge = interaksi"),
        ("4", "User Leaning\n& Community", "Leaning ∈ [-1,+1] + Leiden + Infomap\n(parallel community detection)"),
        ("5", "ECR 2.0 &\nValidasi", "ECR = intra/inter agreement\n+ Bootstrap CI, Permutation Test, Quality Metrics"),
    ]

    colors_hex = [
        RGBColor(0x42, 0xa5, 0xf5),
        RGBColor(0x66, 0xbb, 0x6a),
        RGBColor(0xff, 0xa7, 0x26),
        RGBColor(0xab, 0x47, 0xbc),
        RGBColor(0xef, 0x53, 0x50),
    ]

    box_w = Inches(2.2)
    box_h = Inches(2.0)
    start_x = Inches(0.6)
    gap = Inches(0.35)
    y_pos = Inches(2.2)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors_hex[i]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Tahap {num}"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = ""
        p3.font.size = Pt(6)

        p4 = tf.add_paragraph()
        p4.text = desc
        p4.font.size = Pt(11)
        p4.font.color.rgb = WHITE
        p4.alignment = PP_ALIGN.CENTER

        # Arrow between boxes
        if i < len(steps) - 1:
            ax = x + box_w
            ay = y_pos + box_h / 2
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay - Inches(0.15),
                gap, Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

    # Output description at bottom
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(1.5),
                "Output: 10 visualisasi PNG + CSV (users, communities) + TXT (ECR metrics, comparison) + Verdict echo chamber / non-echo chamber",
                font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 6: Tahap 1 — Sentiment Analysis
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Tahap 1: Preprocessing & Sentiment Analysis",
                   "Fine-tuned XLM-RoBERTa untuk Bahasa Indonesia")

    items_sa = [
        "Model: cardiffnlp/twitter-xlm-roberta-base-sentiment-multilingual",
        "• Di-fine-tune pada dataset sentimen Indonesia",
        "• 3-class: negative, neutral, positive",
        "• Output probabilitas: p_neg, p_neu, p_pos per tweet",
        "",
        "Proses:",
        "• Batch processing (batch_size=64, max_length=128)",
        "• GPU acceleration (CUDA) jika tersedia",
        "• Auto-caching: hasil inferensi disimpan ke CSV",
        "  → tidak perlu re-run saat eksperimen berulang",
        "",
        "Keunggulan vs Model Generic:",
        "• Memahami bahasa informal Indonesia (slang, singkatan)",
        "• Pre-trained pada data Twitter multilingual",
        "• Fine-tuned spesifik untuk konteks isu Indonesia",
        "",
        "Konversi ke User Leaning:",
        "  lean_scalar = p_pos - p_neg  (range [-1, +1])",
        "  -1 = contra, 0 = neutral, +1 = pro topik",
    ]
    add_bullet_slide(slide, items_sa,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 7: Tahap 2 — Graf Interaksi
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Tahap 2: Konstruksi Graf Interaksi",
                   "Directed Weighted Interaction Graph")

    items_graph = [
        "Jenis Interaksi (Edge Types):",
        "  • Reply: user A membalas tweet user B",
        "  • Retweet: user A me-retweet tweet user B",
        "  • Mention: user A menyebut user B di tweet",
        "",
        "Properti Graf:",
        "  • Directed: arah interaksi dipertahankan",
        "  • Weighted: bobot = jumlah interaksi antar user",
        "  • Node = username (bukan user ID)",
        "  • Filter: hanya user yang memiliki konten (tweet)",
        "",
        "Contoh:",
        "  @user_A → reply → @user_B  (weight=3)",
        "  @user_A → retweet → @user_C  (weight=1)",
        "  @user_B → mention → @user_A  (weight=2)",
    ]
    add_bullet_slide(slide, items_graph,
                     Inches(0.5), Inches(1.5), Inches(6), Inches(5.5),
                     font_size=15, color=BLACK)

    # Try adding network visualization
    net_img = find_viz("08_network.png")
    if net_img:
        slide.shapes.add_picture(net_img, Inches(7), Inches(1.8),
                                 width=Inches(5.5))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 8: Tahap 3 — Community Detection
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Tahap 3: User Leaning & Community Detection",
                   "Dual Algorithm: Leiden + Infomap (Parallel)")

    items_cd = [
        "Leiden Algorithm:",
        "  • Optimasi modularitas (perbaikan dari Louvain)",
        "  • Menjamin connected communities",
        "  • Parameter: resolution (mempengaruhi granularity)",
        "  • Lebih cepat dan stabil untuk graf besar",
        "",
        "Infomap Algorithm:",
        "  • Optimasi information flow (random walk compression)",
        "  • Deteksi komunitas berdasarkan arus informasi",
        "  • Efektif untuk graf directed",
        "  • Perspektif berbeda dari modularity-based",
        "",
        "Mengapa Dual?",
        "  • Cross-validation antar algoritma",
        "  • NMI mengukur konsistensi hasil",
        "  • Meningkatkan confidence pada hasil ECR",
        "  • Jika NMI rendah → perlu investigasi lebih lanjut",
    ]
    add_bullet_slide(slide, items_cd,
                     Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5),
                     font_size=15, color=BLACK)

    comm_img = find_viz("03_community_sizes.png")
    if comm_img:
        slide.shapes.add_picture(comm_img, Inches(7.2), Inches(1.8),
                                 width=Inches(5.5))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 9: Tahap 4 — ECR & Validasi
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Tahap 4: ECR 2.0 & Validasi Statistik",
                   "Echo Chamber Ratio + Multiple Validation")

    items_ecr = [
        "Formula ECR 2.0:",
        "                E[intra-community agreement]",
        "  ECR  =  ────────────────────────────────────",
        "                E[inter-community agreement]",
        "",
        "  • intra = avg similarity user dalam komunitas SAMA",
        "  • inter = avg similarity user dari komunitas BERBEDA",
        "  • ECR < threshold → ECHO CHAMBER DETECTED",
        "",
        "Validasi Statistik:",
        "  1. Null Model Threshold — ECR dari N graf random",
        "  2. Bootstrap 95% CI — resample users 100x,",
        "     hitung ECR per sample → interval kepercayaan",
        "  3. Permutation Test — shuffle leaning labels 100x,",
        "     hitung p-value → signifikansi statistik",
        "",
        "Community Quality Metrics:",
        "  • Modularity > 0.3 → meaningful communities",
        "  • Coverage — fraksi intra-community edges",
        "  • Conductance < 0.5 → well-separated",
        "  • NMI (Leiden ↔ Infomap) > 0.5 → konsisten",
    ]
    add_bullet_slide(slide, items_ecr,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 10: Tahap 5 — Analisis Lanjutan
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Tahap 5: Analisis Lanjutan & Visualisasi")

    items_adv = [
        "Homophily Analysis:",
        "  • Pearson correlation: user leaning vs neighbor leaning",
        "  • r > 0.3 → homophily exists (birds of a feather)",
        "",
        "Diffusion Bias (Independent Cascade):",
        "  • Simulasi penyebaran informasi di graf",
        "  • Ukur apakah informasi cenderung menyebar",
        "    ke user yang sepandangan (bias) atau tidak",
        "  • Metrik: slope, R², bias_ratio, Cohen's d",
        "",
        "Community Classification:",
        "  • Setiap komunitas diklasifikasi: pro/neutral/contra",
        "  • Berdasarkan rata-rata leaning anggota komunitas",
        "  • Top users per stance (most influential)",
    ]
    add_bullet_slide(slide, items_adv,
                     Inches(0.5), Inches(1.5), Inches(6), Inches(5.5),
                     font_size=15, color=BLACK)

    # 10 Visualizations list
    viz_items = [
        "10 Visualisasi Otomatis:",
        "",
        "01. Calibration Impact (reliability diagrams)",
        "02. User Leaning (histogram + scatter)",
        "03. Community Sizes (bar chart)",
        "04. ECR Results & Homophily",
        "05. Diffusion Bias (regression + histogram)",
        "06. Pro/Contra Analysis (stance + top users)",
        "07. Sentiment Shifts (trajectory + flow)",
        "08. Network Graph (spring layout)",
        "09. Degree Distribution (total/in/out)",
        "10. Sentiment Sample (10 contoh)",
    ]
    add_bullet_slide(slide, viz_items,
                     Inches(6.8), Inches(1.5), Inches(6), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 11: Implementasi GUI
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Implementasi: Web GUI (NiceGUI)",
                   "Antarmuka interaktif untuk pipeline ECR 2.0")

    items_gui = [
        "Teknologi:",
        "  • Python + NiceGUI (web framework)",
        "  • Berjalan di http://localhost:8080",
        "  • Responsive, real-time progress logging",
        "",
        "Fitur GUI:",
        "  • Input: path dataset CSV",
        "  • Konfigurasi parameter dengan tooltip info (i)",
        "  • Real-time log panel (setiap tahap pipeline)",
        "  • Hasil langsung ditampilkan di browser",
        "  • 10 visualisasi otomatis di halaman web",
        "",
        "Fitur Tambahan:",
        "  • Load Previous Results — untuk demo/presentasi",
        "    (load hasil pipeline sebelumnya secara instan)",
        "  • Auto-cache sentimen (skip inference jika sudah ada)",
        "  • GPU acceleration otomatis jika CUDA tersedia",
        "",
        "File: main_gui.py (~950 baris)",
        "Dependencies: 13 packages (requirements.txt)",
    ]
    add_bullet_slide(slide, items_gui,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 12: Dataset & Eksperimen
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Dataset & Eksperimen")

    items_ds = [
        "Dataset Percakapan Twitter/X Indonesia:",
        "",
        "  Dataset                     Rows         Users      Topik",
        "  ──────────────────────────────────────────────────────────",
        "  Boikot McDonald's          ~515K        ~168K      Boikot produk",
        "  Indonesia Gelap             varies       varies     Isu nasional",
        "  Vaksinasi                   varies       varies     Pro/anti vaksin",
        "",
        "Format Input CSV:",
        "  • from_id, from_username — identitas pengirim",
        "  • cleaned_text — teks tweet yang sudah dibersihkan",
        "  • reply_to_user_id — target reply",
        "  • retweet_from_user_id — sumber retweet",
        "  • mentioned — daftar user yang di-mention",
        "",
        "Eksperimen dilakukan dengan:",
        "  • Leiden resolution: auto-optimized (~5.2 untuk Boikot)",
        "  • Infomap: directed mode, 2-level",
        "  • Bootstrap: 100 samples, Permutation: 100 iterations",
        "  • Diffusion: Independent Cascade, p_activate=0.05",
    ]
    add_bullet_slide(slide, items_ds,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 13: Hasil ECR
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Hasil: ECR — Leiden vs Infomap",
                   "Dataset: Boikot McDonald's (~515K tweets, ~168K users)")

    # Results table as text
    items_res = [
        "Perbandingan Leiden vs Infomap:",
        "",
        "  Metrik                      Leiden      Infomap",
        "  ─────────────────────────────────────────────────",
        "  Jumlah Komunitas              2138         3663",
        "  Intra-community agreement     0.828        0.831",
        "  Inter-community agreement     0.849        0.848",
        "  ECR 2.0 Ratio                 0.976        0.981",
        "  Threshold (Null Model)        1.000        1.000",
        "  Klasifikasi               Non-Echo    Non-Echo",
        "",
        "  Verdict: FENOMENA NON-ECHO CHAMBER",
        "  (ECR > threshold pada kedua algoritma)",
    ]
    add_bullet_slide(slide, items_res,
                     Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5),
                     font_size=15, color=BLACK)

    ecr_img = find_viz("04_ecr_results.png")
    if ecr_img:
        slide.shapes.add_picture(ecr_img, Inches(7), Inches(1.8),
                                 width=Inches(5.8))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 14: Community Quality
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Hasil: Community Quality & Validation Metrics")

    items_cq = [
        "Community Quality Validation:",
        "",
        "  Metrik                      Leiden      Infomap",
        "  ─────────────────────────────────────────────────",
        "  Modularity                   0.634        0.640",
        "  Coverage                     0.619        0.712",
        "  Performance                  0.986        0.924",
        "  Conductance (mean)           0.031        0.052",
        "",
        "  NMI (Leiden ↔ Infomap):  0.679",
        "  Homophily (Pearson r):   0.490",
        "",
        "Interpretasi:",
        "  ✓ Modularity > 0.3 → komunitas meaningful",
        "  ✓ Conductance < 0.5 → komunitas well-separated",
        "  ✓ NMI > 0.5 → Leiden & Infomap konsisten",
        "  ✓ Coverage tinggi → mayoritas edge dalam komunitas",
        "  ✓ Homophily moderat → ada kecenderungan berinteraksi",
        "    dengan user sepandangan, tapi tidak ekstrem",
    ]
    add_bullet_slide(slide, items_cq,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 15: Visualisasi 1 — Sentimen & Leaning
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Hasil Visualisasi (1): Sentimen & User Leaning")

    cal_img = find_viz("01_calibration_impact.png")
    lean_img = find_viz("02_user_leaning.png")

    if cal_img:
        slide.shapes.add_picture(cal_img, Inches(0.3), Inches(1.5),
                                 width=Inches(6.2))
    if lean_img:
        slide.shapes.add_picture(lean_img, Inches(6.8), Inches(1.5),
                                 width=Inches(6.2))

    if not cal_img and not lean_img:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    "(Gambar visualisasi akan ditampilkan saat file tersedia di folder hasil)",
                    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 16: Visualisasi 2 — Komunitas & ECR
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Hasil Visualisasi (2): Komunitas & ECR")

    comm_img = find_viz("03_community_sizes.png")
    ecr_img2 = find_viz("04_ecr_results.png")

    if comm_img:
        slide.shapes.add_picture(comm_img, Inches(0.3), Inches(1.5),
                                 width=Inches(6.2))
    if ecr_img2:
        slide.shapes.add_picture(ecr_img2, Inches(6.8), Inches(1.5),
                                 width=Inches(6.2))

    if not comm_img and not ecr_img2:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    "(Gambar visualisasi akan ditampilkan saat file tersedia di folder hasil)",
                    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 17: Visualisasi 3 — Diffusion & Network
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Hasil Visualisasi (3): Diffusion Bias & Network")

    diff_img = find_viz("05_diffusion_bias.png")
    pro_img = find_viz("06_pro_contra.png")

    if diff_img:
        slide.shapes.add_picture(diff_img, Inches(0.3), Inches(1.5),
                                 width=Inches(6.2))
    if pro_img:
        slide.shapes.add_picture(pro_img, Inches(6.8), Inches(1.5),
                                 width=Inches(6.2))

    if not diff_img and not pro_img:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    "(Gambar visualisasi akan ditampilkan saat file tersedia di folder hasil)",
                    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 18: Analisis Akhir
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Analisis Akhir")

    items_af = [
        "Temuan Utama — Boikot McDonald's:",
        "",
        "1. Tidak ditemukan echo chamber yang kuat (ECR ≈ 0.98)",
        "   • Intra-community agreement ≈ inter-community agreement",
        "   • Pengguna pro dan kontra masih saling berinteraksi",
        "",
        "2. Komunitas memiliki modularity tinggi (>0.63)",
        "   • Pembagian komunitas jelas dan meaningful",
        "   • Bukan artefak dari algoritma deteksi",
        "",
        "3. Homophily moderat (r = 0.49)",
        "   • Ada kecenderungan berinteraksi dengan user sepandangan",
        "   • Namun tidak cukup kuat untuk membentuk echo chamber",
        "",
        "4. Diffusion bias positif tapi lemah",
        "   • Informasi sedikit bias ke user sepandangan (slope=0.30)",
        "   • Cohen's d = -0.15 → efek size sangat kecil",
        "",
        "5. Leiden dan Infomap memberikan hasil konsisten",
        "   • NMI = 0.679 → kedua algoritma menemukan struktur serupa",
        "   • Meningkatkan confidence pada kesimpulan",
    ]
    add_bullet_slide(slide, items_af,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 19: Kesimpulan
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Kesimpulan")

    items_conc = [
        "1. Pipeline ECR 2.0 berhasil dikembangkan sebagai tools otomatis",
        "   end-to-end untuk mengukur echo chamber di media sosial Indonesia.",
        "",
        "2. Fine-tuned XLM-RoBERTa efektif untuk analisis sentimen",
        "   bahasa Indonesia informal (tweet, slang, singkatan).",
        "",
        "3. Pendekatan dual community detection (Leiden + Infomap) memberikan",
        "   cross-validation yang meningkatkan reliabilitas hasil — dibuktikan",
        "   dengan NMI > 0.5 pada semua eksperimen.",
        "",
        "4. Multiple validation metrics (Bootstrap CI, Permutation Test,",
        "   Modularity, Coverage, Conductance) memberikan bukti statistik",
        "   yang kuat mengenai ada/tidaknya echo chamber.",
        "",
        "5. Web GUI (NiceGUI) memungkinkan analisis interaktif tanpa",
        "   memerlukan kemampuan programming — cocok untuk peneliti,",
        "   jurnalis, dan policy-maker.",
        "",
        "6. Pada dataset Boikot McDonald's (~515K tweet), hasil menunjukkan",
        "   fenomena non-echo chamber (ECR > threshold) dengan homophily",
        "   moderat, mengindikasikan dialog lintas pandangan masih terjadi.",
    ]
    add_bullet_slide(slide, items_conc,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 20: Saran & Future Work
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    make_title_bar(slide, "Saran & Future Work")

    items_fw = [
        "Saran untuk Pengembangan Selanjutnya:",
        "",
        "1. Temporal Analysis",
        "   • Analisis evolusi echo chamber dari waktu ke waktu",
        "   • Deteksi kapan echo chamber mulai terbentuk/hilang",
        "",
        "2. Dataset Lebih Beragam",
        "   • Uji pada lebih banyak topik kontroversial Indonesia",
        "   • Perbandingan lintas platform (Twitter, Instagram, TikTok)",
        "",
        "3. Model Sentimen Lebih Lanjut",
        "   • Aspect-Based Sentiment Analysis (ABSA)",
        "   • Integrasi dengan LLM (GPT, Gemini) untuk zero-shot",
        "",
        "4. Skalabilitas",
        "   • Optimasi untuk dataset >1M tweet",
        "   • Distributed computing (Spark, Dask)",
        "",
        "5. Deployment",
        "   • Deploy sebagai web service (cloud-hosted)",
        "   • API endpoint untuk integrasi dengan sistem lain",
        "   • Dashboard real-time monitoring echo chamber",
    ]
    add_bullet_slide(slide, items_fw,
                     Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
                     font_size=15, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 21: Terima Kasih
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
                "Terima Kasih",
                font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
                "Pengukuran Echo Chamber di Media Sosial Indonesia\n"
                "Menggunakan Framework ECR 2.0",
                font_size=20, color=LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                "Amalia Putri — Teknik Informatika — STEI ITB — 2025",
                font_size=16, color=RGBColor(0xbb, 0xde, 0xfb),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
                "Pertanyaan?",
                font_size=24, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Save
    output_path = OUT_DIR / "presentasi_ecr2.pptx"
    prs.save(str(output_path))
    print(f"[OK] Presentation saved: {output_path}")


if __name__ == "__main__":
    build_pptx()
